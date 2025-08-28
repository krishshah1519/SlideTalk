import io
import base64
import magic
from typing import List, Dict, Any
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture


def _table_to_markdown(shape: GraphicFrame) -> str:
    """Converts a table shape to a Markdown string."""

    table = shape.table
    markdown_table = ""
    try:
        header_cells = [cell.text for cell in table.rows[0].cells]
        markdown_table += "| " + " | ".join(header_cells) + " |\n"
        markdown_table += "| " + " | ".join(["---"] * len(header_cells)) + " |\n"
        for row in list(table.rows)[1:]:
            data_cells = [cell.text for cell in row.cells]
            markdown_table += "| " + " | ".join(data_cells) + " |\n"
    except IndexError:
        return "Table is empty or improperly formatted."
    return markdown_table


def _image_to_base64(shape: Picture) -> str:
    """Converts a picture shape to a Base64 data URI."""

    image_blob = shape.image.blob
    mime_type = magic.from_buffer(image_blob, mime=True)
    base64_str = base64.b64encode(image_blob).decode('utf-8')
    return f"data:{mime_type};base64,{base64_str}"


def extract_content_from_ppt(ppt_file: io.BytesIO) -> List[Dict[str, Any]]:
    """Extracts all content (text, tables, images) from a PPTX file stream."""

    prs = pptx.Presentation(ppt_file)
    slides_content = []

    for i, slide in enumerate(prs.slides):
        slide_elements = []
        for shape in sorted(slide.shapes, key=lambda s: (s.top, s.left)):
            if shape.has_text_frame and shape.text.strip():
                slide_elements.append({"type": "text", "data": shape.text})
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                slide_elements.append({"type": "table", "data": _table_to_markdown(shape)})
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_elements.append({"type": "image", "data": _image_to_base64(shape)})


        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text

        slides_content.append({
            "slide_number": i + 1,
            "content": slide_elements,
            "notes": notes,
        })
    return slides_content